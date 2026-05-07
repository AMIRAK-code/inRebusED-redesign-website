"""One-shot inspector: dump text + image sizes for every slide."""
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = r"C:\Users\faraz\Desktop\inrebus-website\.claude\inRebus_DL_2026_EN.pptx"

prs = Presentation(PPTX_PATH)
total = len(prs.slides)
print(f"Total slides: {total}\n")

for idx, slide in enumerate(prs.slides, 1):
    imgs = []
    texts = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            imgs.append(f"  IMG [{shape.width//914400:.2f}\" x {shape.height//914400:.2f}\"] {img.content_type} {len(img.blob)//1024}KB name={shape.name}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(f"  TEXT: {t[:120]}")
    print(f"=== Slide {idx} ===")
    for t in texts:
        print(t)
    for i in imgs:
        print(i)
    print()
