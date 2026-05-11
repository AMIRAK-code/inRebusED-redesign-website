"""Dump full text from case-study slides (no truncation)."""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = Path(__file__).parent.parent / ".claude" / "inRebus_DL_2026_EN.pptx"
prs = Presentation(PPTX_PATH)
slides = prs.slides

CONTENT_SLIDES = [9,10,11,12,13,15,16,17,19,20,21,22,24,25,26,27,28,29,30]

for idx in CONTENT_SLIDES:
    slide = slides[idx]
    print("=" * 60)
    print("SLIDE {} (0-idx {})".format(idx + 1, idx))
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print("  | " + t)
    print()
