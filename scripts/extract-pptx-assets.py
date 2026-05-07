"""
Extract assets from inRebus_DL_2026_EN.pptx.

Outputs:
  public/case-studies/<slug>.jpg  - largest content image per case-study slide (<= 1600px wide, 85% JPEG)
  public/clients/<N>.png          - each logo from slide 32 (trim transparent padding)
  public/og.jpg                   - 1200x630 OG card
  public/favicon.svg              - "iR" mark on navy background

Run from repo root:
    python scripts/extract-pptx-assets.py
"""

import io
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw, ImageFont
except ImportError as e:
    sys.exit("Missing dependency: {}\nRun: python -m pip install python-pptx Pillow".format(e))

REPO       = Path(__file__).parent.parent
PPTX_PATH  = REPO / ".claude" / "inRebus_DL_2026_EN.pptx"
CASE_DIR   = REPO / "public" / "case-studies"
CLIENT_DIR = REPO / "public" / "clients"

CASE_DIR.mkdir(parents=True, exist_ok=True)
CLIENT_DIR.mkdir(parents=True, exist_ok=True)

# (0-based slide index, output slug)
# Section-header slides (0-idx 8,14,18,23) are deliberately skipped.
CASE_STUDY_SLIDES = [
    ( 9,  "supply-chain-academy"),           # slide 10
    (10,  "digital-liguria"),                # slide 11
    (11,  "dussmann-academy"),               # slide 12
    (12,  "technology-and-digital-academy"), # slide 13 - IVECO internal
    (13,  "marelli-campus"),                 # slide 14
    (15,  "esselunga-onboarding"),           # slide 16
    (16,  "onboarding-dussmann"),            # slide 17
    (17,  "maserati-history"),               # slide 18
    (19,  "m-346-c-27j-familiarisation"),    # slide 20
    (20,  "clienteling-training"),           # slide 21
    (21,  "gran-turismo"),                   # slide 22
    (22,  "iveco-edaily"),                   # slide 23
    (24,  "corporate-brand-identity"),       # slide 25
    (25,  "zero-tolerance"),                 # slide 26
    (26,  "genuine-parts-training"),         # slide 27
    (27,  "anti-money-laundering"),          # slide 28
    (28,  "hygiene-and-nutrition"),          # slide 29
    (29,  "corporate-liability-231"),        # slide 30
    (30,  "crm-dealers"),                    # slide 31
]

CLIENTS_SLIDE_IDX = 31   # slide 32, 0-based

MAX_WIDTH    = 1600
JPEG_QUALITY = 85
LOGO_MAX_BYTES = 60_000  # blobs under this threshold are treated as decorative


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def collect_images(shape):
    """Recursively yield (blob, content_type) for all PICTURE shapes."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            yield img.blob, img.content_type
        except ValueError:
            pass  # linked (not embedded) image — skip
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from collect_images(child)


def largest_content_image(slide):
    """Return blob of the biggest image on a slide, ignoring small logos."""
    candidates = [
        blob
        for shape in slide.shapes
        for blob, _ in collect_images(shape)
        if len(blob) > LOGO_MAX_BYTES
    ]
    return max(candidates, key=len) if candidates else None


def save_as_jpeg(blob, dest):
    img = Image.open(io.BytesIO(blob)).convert("RGB")
    if img.width > MAX_WIDTH:
        ratio = MAX_WIDTH / img.width
        img = img.resize((MAX_WIDTH, round(img.height * ratio)), Image.LANCZOS)
    img.save(dest, "JPEG", quality=JPEG_QUALITY, optimize=True)
    print("  + {}  ({}x{})".format(dest.name, img.width, img.height))


def trim_and_save_png(blob, dest):
    img = Image.open(io.BytesIO(blob)).convert("RGBA")
    bbox = img.getbbox()
    if bbox:
        img = img.crop(bbox)
    img.save(dest, "PNG")
    print("  + {}  ({}x{})".format(dest.name, img.width, img.height))


def make_placeholder(dest):
    """Navy placeholder with orange left bar."""
    img  = Image.new("RGB", (1200, 675), (10, 22, 40))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, 8, 675], fill=(245, 130, 32))
    img.save(dest, "JPEG", quality=85)
    print("  ~ placeholder {}".format(dest.name))


def load_font(size):
    for name in ["arialbd.ttf", "Arial Bold.ttf", "DejaVuSans-Bold.ttf", "arial.ttf"]:
        try:
            return ImageFont.truetype(name, size)
        except Exception:
            pass
    return ImageFont.load_default()


# ---------------------------------------------------------------------------
# OG image
# ---------------------------------------------------------------------------

def make_og_image():
    W, H   = 1200, 630
    NAVY   = (10, 22, 40)
    ORANGE = (245, 130, 32)
    WHITE  = (247, 247, 247)
    MUTED  = (180, 180, 200)

    img  = Image.new("RGB", (W, H), NAVY)
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, 8, H], fill=ORANGE)

    f72 = load_font(72)
    f40 = load_font(40)
    f28 = load_font(28)

    draw.text((72,  80), "inRebus Digital Learning",       font=f28, fill=ORANGE)
    draw.text((72, 160), "CREATIVE",                        font=f72, fill=WHITE)
    draw.text((72, 250), "OUTSIDE THE BOX",                 font=f72, fill=ORANGE)
    draw.text((72, 380), "Custom e-learning, gamification &", font=f40, fill=MUTED)
    draw.text((72, 430), "digital academies since 2003",    font=f40, fill=MUTED)
    draw.text((72, 560), "inrebus.education",               font=f28, fill=ORANGE)

    dest = REPO / "public" / "og.jpg"
    img.save(dest, "JPEG", quality=90)
    print("  + og.jpg  ({}x{})".format(W, H))


# ---------------------------------------------------------------------------
# Favicon SVG
# ---------------------------------------------------------------------------

def make_favicon():
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 32 32">\n'
        '  <rect width="32" height="32" rx="6" fill="#0A1628"/>\n'
        '  <text x="4" y="24" font-family="Arial,sans-serif" font-weight="bold"\n'
        '        font-size="18" fill="#F58220">iR</text>\n'
        '</svg>\n'
    )
    dest = REPO / "public" / "favicon.svg"
    dest.write_text(svg, encoding="utf-8")
    print("  + favicon.svg")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    if not PPTX_PATH.exists():
        sys.exit("PPTX not found: {}".format(PPTX_PATH))

    print("Loading {} ...".format(PPTX_PATH.name))
    prs    = Presentation(PPTX_PATH)
    slides = prs.slides

    # -- Case-study images ---------------------------------------------------
    print("\n-- Case-study images ({} slides) --".format(len(CASE_STUDY_SLIDES)))
    for slide_idx, slug in CASE_STUDY_SLIDES:
        dest = CASE_DIR / "{}.jpg".format(slug)
        if dest.exists():
            print("  skip (exists) {}".format(dest.name))
            continue
        blob = largest_content_image(slides[slide_idx])
        if blob is None:
            print("  ! no image on slide {} ({})".format(slide_idx + 1, slug))
            make_placeholder(dest)
        else:
            save_as_jpeg(blob, dest)

    # -- Client logos --------------------------------------------------------
    print("\n-- Client logos (slide 32) --")
    blobs = [
        blob
        for shape in slides[CLIENTS_SLIDE_IDX].shapes
        for blob, _ in collect_images(shape)
    ]
    if blobs:
        for i, blob in enumerate(blobs, 1):
            trim_and_save_png(blob, CLIENT_DIR / "{:02d}.png".format(i))
    else:
        print("  ! No image shapes found on slide 32 (logos may be vector shapes)")
        print("    public/clients/ is empty - clients will render as text chips")

    # -- OG image ------------------------------------------------------------
    print("\n-- OG image --")
    make_og_image()

    # -- Favicon -------------------------------------------------------------
    print("\n-- Favicon --")
    make_favicon()

    print("\nDone.")


if __name__ == "__main__":
    main()
