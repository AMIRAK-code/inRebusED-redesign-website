"""
Map extracted client logos (public/clients/01.png…40.png) to company names.

Walks slide 32 in the exact same insertion order used by extract-pptx-assets.py,
then attempts to match each logo to a nearby text label on the same slide.

Outputs:
  scripts/client-logo-map.json  — { "01": "Esselunga", "02": "IVECO", … }

Run from repo root:
    python scripts/map-client-logos.py
"""

import hashlib
import json
import sys
import io
from pathlib import Path

# Force UTF-8 output on Windows
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("Missing dependency. Run: python -m pip install python-pptx")

REPO      = Path(__file__).parent.parent
PPTX_PATH = REPO / ".claude" / "inRebus_DL_2026_EN.pptx"
OUT_JSON  = Path(__file__).parent / "client-logo-map.json"

SLIDE_IDX    = 31   # slide 32, 0-based (same as extract script)
PROXIMITY_PT = 80   # max vertical distance (pt) to associate a text label


# ── Unit helpers ─────────────────────────────────────────────────────────────

def emu_to_pt(emu: int) -> float:
    return emu / 12700


def shape_bbox(shape):
    return {
        'l': emu_to_pt(shape.left),
        't': emu_to_pt(shape.top),
        'r': emu_to_pt(shape.left + shape.width),
        'b': emu_to_pt(shape.top + shape.height),
    }


# ── Image collection (matches extract-pptx-assets.py order exactly) ──────────

def collect_images(shape):
    """Recursively yield (blob, content_type) — same logic as extract script."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            yield img.blob, img.content_type, shape.name, shape_bbox(shape)
        except ValueError:
            pass
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from collect_images(child)


# ── Text collection ───────────────────────────────────────────────────────────

def collect_texts(slide):
    """Return list of {text, bbox} for every non-empty text frame on the slide."""
    results = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        lines = [
            para.text.strip()
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        ]
        text = ' '.join(lines)
        if text:
            try:
                results.append({'text': text, 'bbox': shape_bbox(shape)})
            except Exception:
                pass
    return results


# ── Spatial matching ──────────────────────────────────────────────────────────

def nearest_text(img_bbox, texts, threshold_pt=PROXIMITY_PT):
    """
    Find the text shape whose vertical centre is within threshold_pt of the
    image's vertical centre AND whose horizontal extent overlaps the image.
    Returns the text string, or None.
    """
    img_cy = (img_bbox['t'] + img_bbox['b']) / 2

    best_text = None
    best_dist = threshold_pt

    for t in texts:
        tb = t['bbox']
        # Require at least some horizontal overlap (with 30pt slack)
        if tb['r'] < img_bbox['l'] - 30 or tb['l'] > img_bbox['r'] + 30:
            continue
        t_cy = (tb['t'] + tb['b']) / 2
        dist = abs(t_cy - img_cy)
        if dist < best_dist:
            best_dist = dist
            best_text = t['text']

    return best_text


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if not PPTX_PATH.exists():
        sys.exit(f"PPTX not found: {PPTX_PATH}")

    prs   = Presentation(PPTX_PATH)
    slide = prs.slides[SLIDE_IDX]

    # Collect images in the same iteration order as extract-pptx-assets.py
    image_entries = []
    for shape in slide.shapes:
        for blob, ctype, sname, bbox in collect_images(shape):
            image_entries.append({
                'blob':       blob,
                'shape_name': sname,
                'bbox':       bbox,
                'hash':       hashlib.md5(blob).hexdigest()[:8],
            })

    # Collect text labels on the slide
    texts = collect_texts(slide)

    print(f"Slide 32: {len(image_entries)} images, {len(texts)} text shapes\n")
    print("-- All text shapes on slide 32 -------------------------------------")
    for t in texts:
        bb = t['bbox']
        print(f"  [{bb['l']:6.0f},{bb['t']:6.0f}]  {t['text'][:80]}")
    print()

    # Build mapping
    mapping: dict[str, str | None] = {}
    unknown: list[tuple[str, str]] = []

    print("-- Logo mapping ----------------------------------------------------")
    for i, entry in enumerate(image_entries, 1):
        key = f"{i:02d}"

        # 1. Nearest text label by spatial proximity
        name = nearest_text(entry['bbox'], texts)

        # 2. Fall back to shape name if it looks intentional
        if name is None:
            sname = entry['shape_name']
            if sname and not sname.lower().startswith(('picture', 'image', 'group')):
                name = sname

        mapping[key] = name

        if name:
            print(f"  {key} → {name}")
        else:
            bb = entry['bbox']
            label = f"UNKNOWN  (hash:{entry['hash']}  pos:[{bb['l']:.0f},{bb['t']:.0f}])"
            print(f"  {key} → {label}")
            unknown.append((key, entry['hash']))

    if unknown:
        print(f"\nWARN: {len(unknown)} unmapped logo(s) - identify by hash + position above:")
        for key, h in unknown:
            print(f"   {key}  md5[:8]={h}  -> public/clients/{key}.png")

    OUT_JSON.write_text(
        json.dumps(mapping, indent=2, ensure_ascii=False),
        encoding='utf-8',
    )
    print(f"\nSaved: {OUT_JSON}")


if __name__ == "__main__":
    main()
